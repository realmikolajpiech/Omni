import os
import csv
import io
import logging

# Setup logging
logger = logging.getLogger(__name__)

IMAGE_EXTS = {".jpg", ".jpeg", ".png", ".webp", ".bmp", ".gif"}

TEXT_EXTENSIONS = {
    '.txt', '.md', '.py', '.js', '.html', '.css', '.c', '.cpp', '.h', '.sh',
    '.json', '.yml', '.yaml', '.toml', '.ini', '.cfg', '.conf', '.java', '.rs', '.go',
    '.pdf', '.docx', '.rtf', '.csv', '.xlsx', '.xls', '.pptx', '',
}


def is_text_file(path):
    _, ext = os.path.splitext(path)
    return ext.lower() in TEXT_EXTENSIONS


def is_image_file(path):
    _, ext = os.path.splitext(path)
    return ext.lower() in IMAGE_EXTS


# ---------------------------------------------------------------------------
# Per-format text extraction
# ---------------------------------------------------------------------------

def read_pdf(path):
    """Extracts text from a PDF file."""
    try:
        import pypdf
        text = ""
        with open(path, 'rb') as f:
            reader = pypdf.PdfReader(f)
            for page in reader.pages:
                extracted = page.extract_text()
                if extracted:
                    text += extracted + "\n"
        return text
    except Exception as e:
        logger.error(f"Error reading PDF {path}: {e}")
        return ""


def read_docx(path):
    """Extracts text from a DOCX file."""
    try:
        import docx
        doc = docx.Document(path)
        text = "\n".join([para.text for para in doc.paragraphs])
        return text
    except Exception as e:
        logger.error(f"Error reading DOCX {path}: {e}")
        return ""


def read_xlsx(path):
    """Extracts cell text from an Excel workbook (.xlsx / .xls)."""
    try:
        import openpyxl
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts = []
        for ws in wb.worksheets:
            for row in ws.iter_rows(values_only=True):
                cells = [str(c) for c in row if c is not None]
                if cells:
                    parts.append(" | ".join(cells))
        wb.close()
        return "\n".join(parts)
    except Exception as e:
        logger.error(f"Error reading XLSX {path}: {e}")
        return ""


def read_csv_file(path):
    """Extracts text from a CSV file."""
    try:
        with open(path, 'r', encoding='utf-8', errors='ignore') as f:
            sample = f.read(8192)
            f.seek(0)
            try:
                dialect = csv.Sniffer().sniff(sample)
            except csv.Error:
                dialect = csv.excel
            reader = csv.reader(f, dialect)
            parts = []
            for row in reader:
                cells = [c for c in row if c.strip()]
                if cells:
                    parts.append(" | ".join(cells))
            return "\n".join(parts)
    except Exception as e:
        logger.error(f"Error reading CSV {path}: {e}")
        return ""


def read_pptx(path):
    """Extracts slide text from a PowerPoint file."""
    try:
        from pptx import Presentation
        prs = Presentation(path)
        parts = []
        for slide in prs.slides:
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        logger.error(f"Error reading PPTX {path}: {e}")
        return ""


def strip_rtf(rtf_text):
    """Simple RTF stripper using regex."""
    import re
    text = rtf_text
    text = re.sub(r'\\[a-zA-Z]+(-?[0-9]*) ?', ' ', text)
    text = text.replace('{', '').replace('}', '')
    text = re.sub(r'\\', ' ', text)
    text = re.sub(r'\s+', ' ', text).strip()
    return text


# ---------------------------------------------------------------------------
# Unified content processing
# ---------------------------------------------------------------------------

_CONTENT_READERS = {
    '.pdf': read_pdf,
    '.docx': read_docx,
    '.xlsx': read_xlsx,
    '.xls': read_xlsx,
    '.csv': read_csv_file,
    '.pptx': read_pptx,
}


def process_file_content(path, chunk_size=1000, overlap=100):
    """Reads file content and returns chunks."""
    try:
        if os.path.getsize(path) > 5 * 1024 * 1024:
            return []

        _, ext = os.path.splitext(path)
        ext = ext.lower()

        reader = _CONTENT_READERS.get(ext)
        if reader:
            content = reader(path)
        elif ext == '.rtf':
            try:
                with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                    content = strip_rtf(f.read())
            except Exception:
                content = ""
        else:
            with open(path, 'r', encoding='utf-8', errors='ignore') as f:
                content = f.read()

        if not content:
            return []

        chunks = []
        for i in range(0, len(content), chunk_size - overlap):
            chunk = content[i:i + chunk_size]
            if len(chunk) < 50:
                continue
            chunks.append(chunk)
        return chunks
    except Exception as e:
        logger.error(f"Error processing file {path}: {e}")
        return []
